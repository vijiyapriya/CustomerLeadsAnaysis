"""
PowerPoint Presentation Generator
Creates a comprehensive presentation from the Excel data analysis
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pandas as pd
from datetime import datetime
import os

print("=" * 70)
print("Creating PowerPoint Presentation")
print("=" * 70)

# Load data
file_path = r"reports/Raw_File_LS_Updated_Regions_Final.xlsx"
print("\nLoading data...")
df = pd.read_excel(file_path)
print(f"✓ Loaded {len(df):,} rows")

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
TITLE_COLOR = RGBColor(31, 78, 121)  # Dark blue
ACCENT_COLOR = RGBColor(68, 114, 196)  # Blue
TEXT_COLOR = RGBColor(51, 51, 51)  # Dark gray

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 248, 255)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = TITLE_COLOR
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = ACCENT_COLOR
    
    # Date
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = datetime.now().strftime('%B %d, %Y')
    date_para = date_frame.paragraphs[0]
    date_para.alignment = PP_ALIGN.CENTER
    date_para.font.size = Pt(16)
    date_para.font.color.rgb = TEXT_COLOR

def add_content_slide(prs, title, content_list, image_path=None):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = TITLE_COLOR
    
    # Divider line
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1), Inches(9), Inches(0))
    line.line.color.rgb = ACCENT_COLOR
    line.line.width = Pt(3)
    
    # Content
    if image_path and os.path.exists(image_path):
        # Text on left, image on right
        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.5), Inches(5.5))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        for item in content_list:
            p = text_frame.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_COLOR
            p.space_before = Pt(6)
        
        # Add image
        slide.shapes.add_picture(image_path, Inches(5.2), Inches(1.5), width=Inches(4.3))
    else:
        # Full width text
        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        for item in content_list:
            p = text_frame.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(16)
            p.font.color.rgb = TEXT_COLOR
            p.space_before = Pt(8)

def add_table_slide(prs, title, df_data, columns):
    """Add a slide with a table"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = TITLE_COLOR
    
    # Divider line
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1), Inches(9), Inches(0))
    line.line.color.rgb = ACCENT_COLOR
    line.line.width = Pt(3)
    
    # Table
    rows = min(len(df_data), 12) + 1  # +1 for header, max 12 rows
    cols = len(columns)
    
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.4), Inches(9), Inches(5.5)).table
    
    # Header row
    for col_idx, col_name in enumerate(columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_COLOR
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Data rows
    for row_idx in range(min(len(df_data), 12)):
        for col_idx, col_name in enumerate(columns):
            cell = table.cell(row_idx + 1, col_idx)
            value = df_data.iloc[row_idx][col_name]
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            cell.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR

def add_image_slide(prs, title, image_path):
    """Add a slide with a large image"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = TITLE_COLOR
    
    # Divider line
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1), Inches(9), Inches(0))
    line.line.color.rgb = ACCENT_COLOR
    line.line.width = Pt(3)
    
    # Add image
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(0.8), Inches(1.5), width=Inches(8.4))

# Slide 1: Title Slide
print("\n📄 Creating slides...")
print("  1. Title slide")
add_title_slide(prs, 
                "Excel Data Analysis Report",
                "Comprehensive Analysis of Lead Database")

# Slide 2: Executive Summary
print("  2. Executive Summary")
total_leads = len(df)
active_leads = len(df[~df['Lead Stage'].isin(['Disqualified', 'Lost', 'Won', 'Closure - Customer', 'Closure'])])
countries = df['Country'].nunique()
me_region = (df['Region Specific'] == 'ME').sum()
eu_region = (df['Region Specific'] == 'EU').sum()

summary_content = [
    f"📊 Total Records: {total_leads:,}",
    f"✅ Active Leads: {active_leads:,} ({active_leads/total_leads*100:.1f}%)",
    f"🌍 Countries: {countries}",
    f"🏢 Companies: {df['Company Name'].nunique():,}",
    "",
    "Regional Distribution:",
    f"  • ME Region: {me_region:,} ({me_region/total_leads*100:.1f}%)",
    f"  • EU Region: {eu_region:,} ({eu_region/total_leads*100:.1f}%)",
    f"  • USA: {(df['Region Specific']=='USA').sum():,}",
    f"  • Others: {(df['Region Specific']=='Others').sum():,}"
]
add_content_slide(prs, "Executive Summary", summary_content)

# Slide 3: Dataset Overview
print("  3. Dataset Overview")
overview_content = [
    f"Total Rows: {len(df):,}",
    f"Total Columns: {len(df.columns)}",
    f"Data Quality: {(1 - df.isnull().sum().sum()/(len(df)*len(df.columns)))*100:.1f}% complete",
    f"Duplicate Rows: {df.duplicated().sum()} (0%)",
    "",
    "Key Metrics:",
    f"  • Contacts: {(df['Lead Stage']=='Contacts').sum():,}",
    f"  • Leads: {(df['Lead Stage']=='Leads').sum():,}",
    f"  • Disqualified: {(df['Lead Stage']=='Disqualified').sum():,}",
    f"  • Won: {(df['Lead Stage']=='Won').sum():,}",
    f"  • Lost: {(df['Lead Stage']=='Lost').sum():,}"
]
add_content_slide(prs, "Dataset Overview", overview_content)

# Slide 4: Active Leads by Country
print("  4. Active Leads by Country")
if os.path.exists('reports/active_leads_by_country.png'):
    add_image_slide(prs, "Active Leads by Country (Top 15)", 'reports/active_leads_by_country.png')

# Slide 5: Active Leads by Stage - with better layout
print("  5. Active Leads by Stage")
active_leads_df = df[~df['Lead Stage'].isin(['Disqualified', 'Lost', 'Won', 'Closure - Customer', 'Closure'])]
stage_counts = active_leads_df['Lead Stage'].value_counts()

slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "Active Leads Distribution by Stage"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(32)
title_para.font.bold = True
title_para.font.color.rgb = TITLE_COLOR

# Divider line
line = slide.shapes.add_shape(1, Inches(0.5), Inches(1), Inches(9), Inches(0))
line.line.color.rgb = ACCENT_COLOR
line.line.width = Pt(3)

# Add image on left
if os.path.exists('reports/active_leads_by_stage.png'):
    slide.shapes.add_picture('reports/active_leads_by_stage.png', Inches(0.5), Inches(1.5), width=Inches(5.5))

# Add data table on right
text_box = slide.shapes.add_textbox(Inches(6.2), Inches(1.5), Inches(3.3), Inches(5.5))
text_frame = text_box.text_frame
text_frame.word_wrap = True

p = text_frame.add_paragraph()
p.text = "Stage Breakdown:"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = TITLE_COLOR
p.space_after = Pt(10)

for stage, count in stage_counts.items():
    pct = (count / len(active_leads_df)) * 100
    p = text_frame.add_paragraph()
    p.text = f"{str(stage)[:20]}"
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_COLOR
    p.space_after = Pt(2)
    
    p = text_frame.add_paragraph()
    p.text = f"  {count:,} ({pct:.1f}%)"
    p.font.size = Pt(11)
    p.font.color.rgb = ACCENT_COLOR
    p.space_after = Pt(8)

# Slide 6: Top Countries Table
print("  6. Top Countries")
country_counts = df['Country'].value_counts().head(15).reset_index()
country_counts.columns = ['Country', 'Count']
country_counts['Percentage'] = (country_counts['Count'] / len(df) * 100).round(2)
country_counts['Percentage'] = country_counts['Percentage'].astype(str) + '%'
add_table_slide(prs, "Top 15 Countries", country_counts, ['Country', 'Count', 'Percentage'])

# Slide 7: Email Bounced Analysis
print("  7. Email Bounced Analysis")
bounced_mask = df['Last Activity'].astype(str).str.contains('bounce', case=False, na=False)
bounced_count = bounced_mask.sum()
bounced_content = [
    f"Total Email Bounced: {bounced_count:,}",
    f"Percentage of Total: {bounced_count/len(df)*100:.2f}%",
    "",
    "Top 5 Countries with Bounced Emails:",
]

if bounced_count > 0:
    bounced_countries = df[bounced_mask]['Country'].value_counts().head(5)
    for country, count in bounced_countries.items():
        bounced_content.append(f"  • {country}: {count:,} ({count/bounced_count*100:.1f}%)")

add_content_slide(prs, "Email Bounced Analysis", bounced_content,
                 'reports/bounced_by_country.png' if os.path.exists('reports/bounced_by_country.png') else None)

# Slide 8: Industry Distribution
print("  8. Industry Distribution")
if 'Industry Vertical' in df.columns:
    industry_counts = df['Industry Vertical'].value_counts().head(12).reset_index()
    industry_counts.columns = ['Industry Vertical', 'Count']
    industry_counts['Percentage'] = (industry_counts['Count'] / df['Industry Vertical'].notna().sum() * 100).round(2)
    industry_counts['Percentage'] = industry_counts['Percentage'].astype(str) + '%'
    add_table_slide(prs, "Top Industries", industry_counts, ['Industry Vertical', 'Count', 'Percentage'])

# Slide 9: Lead Source Analysis
print("  9. Lead Source Analysis")
if 'Lead Source' in df.columns:
    source_counts = df['Lead Source'].value_counts().head(10).reset_index()
    source_counts.columns = ['Lead Source', 'Count']
    source_counts['Percentage'] = (source_counts['Count'] / df['Lead Source'].notna().sum() * 100).round(2)
    source_counts['Percentage'] = source_counts['Percentage'].astype(str) + '%'
    add_table_slide(prs, "Lead Sources", source_counts, ['Lead Source', 'Count', 'Percentage'])

# Slide 10: Role Analysis - Overview
print("  10. Role Analysis Overview")
role_content = [
    "Key Role Categories Identified:",
    "",
    "• Finance Leads: 30,154 records",
    "• IT Leads: 29,586 records",
    "• CEO: 8,091 records",
    "• CFO: 7,039 records",
    "• HR Leads: 6,741 records",
    "",
    f"Total: 81,611 records across all categories"
]
add_content_slide(prs, "Role Analysis - Overview", role_content,
                 'reports/role_category_totals.png' if os.path.exists('reports/role_category_totals.png') else None)

# Slide 11: Combined Role Analysis by Country (Combining slides 11-15)
print("  11. Combined Role Analysis by Country")
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "Key Roles by Country - Summary"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(32)
title_para.font.bold = True
title_para.font.color.rgb = TITLE_COLOR

# Divider line
line = slide.shapes.add_shape(1, Inches(0.5), Inches(1), Inches(9), Inches(0))
line.line.color.rgb = ACCENT_COLOR
line.line.width = Pt(3)

# Left column - HR, IT, Finance
left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.5), Inches(5.5))
left_frame = left_box.text_frame
left_frame.word_wrap = True

p = left_frame.add_paragraph()
p.text = "HR Leads (6,741) - Top 3:"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = TITLE_COLOR
p.space_after = Pt(4)

for line_text in ["  UAE: 2,088 (31%)", "  Saudi: 1,353 (20%)", "  UK: 733 (11%)"]:
    p = left_frame.add_paragraph()
    p.text = line_text
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_COLOR
    p.space_after = Pt(2)

p = left_frame.add_paragraph()
p.text = ""
p.space_after = Pt(8)

p = left_frame.add_paragraph()
p.text = "IT Leads (29,586) - Top 3:"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = TITLE_COLOR
p.space_after = Pt(4)

for line_text in ["  USA: 4,806 (16%)", "  Germany: 3,852 (13%)", "  UK: 3,696 (12%)"]:
    p = left_frame.add_paragraph()
    p.text = line_text
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_COLOR
    p.space_after = Pt(2)

p = left_frame.add_paragraph()
p.text = ""
p.space_after = Pt(8)

p = left_frame.add_paragraph()
p.text = "Finance Leads (30,154) - Top 3:"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = TITLE_COLOR
p.space_after = Pt(4)

for line_text in ["  UAE: 4,916 (16%)", "  UK: 3,844 (13%)", "  India: 3,394 (11%)"]:
    p = left_frame.add_paragraph()
    p.text = line_text
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_COLOR
    p.space_after = Pt(2)

# Right column - CEO, CFO, Summary
right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(5.5))
right_frame = right_box.text_frame
right_frame.word_wrap = True

p = right_frame.add_paragraph()
p.text = "CEO (8,091) - Top 3:"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = TITLE_COLOR
p.space_after = Pt(4)

for line_text in ["  UK: 1,687 (21%)", "  UAE: 1,178 (15%)", "  USA: 1,092 (14%)"]:
    p = right_frame.add_paragraph()
    p.text = line_text
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_COLOR
    p.space_after = Pt(2)

p = right_frame.add_paragraph()
p.text = ""
p.space_after = Pt(8)

p = right_frame.add_paragraph()
p.text = "CFO (7,039) - Top 3:"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = TITLE_COLOR
p.space_after = Pt(4)

for line_text in ["  USA: 1,175 (17%)", "  Germany: 868 (12%)", "  India: 805 (11%)"]:
    p = right_frame.add_paragraph()
    p.text = line_text
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_COLOR
    p.space_after = Pt(2)

p = right_frame.add_paragraph()
p.text = ""
p.space_after = Pt(12)

# Summary box
p = right_frame.add_paragraph()
p.text = "Key Insights:"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = ACCENT_COLOR
p.space_after = Pt(4)

for line_text in ["• UAE leads in HR & Finance", "• USA strong in IT & CFO", "• UK dominant in CEO roles", "• 81,611 total decision makers"]:
    p = right_frame.add_paragraph()
    p.text = line_text
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_COLOR
    p.space_after = Pt(2)

# Slide 12: Regional Updates
print("  12. Regional Updates")
regional_content = [
    "Region Specific Updates:",
    "",
    "ME Region Countries:",
    "  • Saudi Arabia, UAE, Qatar, Kuwait, Oman, Bahrain",
    f"  • Total: {me_region:,} records",
    "",
    "EU Region Countries:",
    "  • UK, Germany, Switzerland, Austria, Belgium,",
    "    Netherlands, Luxembourg, Denmark, Sweden,",
    "    Norway, Finland",
    f"  • Total: {eu_region:,} records"
]
add_content_slide(prs, "Regional Classification", regional_content)

# Slide 13: Data Quality Insights
print("  13. Data Quality")
missing_data = df.isnull().sum().sort_values(ascending=False).head(10)
quality_content = [
    "Top Columns with Missing Data:",
    ""
]
for col, count in missing_data.items():
    pct = count / len(df) * 100
    if count > 0:
        quality_content.append(f"  • {col}: {count:,} ({pct:.1f}%)")

add_content_slide(prs, "Data Quality Insights", quality_content)

# Slide 14: Key Insights
print("  14. Key Insights")
insights_content = [
    "🔍 Key Findings:",
    "",
    f"1. Active Lead Rate: {active_leads/total_leads*100:.1f}% of database",
    "",
    "2. Geographic Focus:",
    "   • UAE, UK, and USA are top markets",
    "   • Strong presence in ME and EU regions",
    "",
    "3. Decision Makers:",
    "   • 81,611 identified in key roles",
    "   • Finance and IT leads dominate",
    "",
    "4. Email Engagement:",
    f"   • {bounced_count:,} bounced emails identified",
    "   • UK and Germany highest bounce rates"
]
add_content_slide(prs, "Key Insights & Findings", insights_content)

# Slide 15: Recommendations
print("  15. Recommendations")
recommendations_content = [
    "📋 Recommendations:",
    "",
    "1. Focus on Active Leads (183,565 records)",
    "",
    "2. Prioritize UAE, UK, and USA markets",
    "",
    "3. Target Finance and IT decision makers",
    "",
    "4. Clean up bounced email addresses",
    "",
    "5. Fill missing data gaps in:",
    "   • Job Titles, Contact Numbers, Practice",
    "",
    "6. Leverage strong ME and EU presence"
]
add_content_slide(prs, "Recommendations", recommendations_content)

# Slide 16: Thank You
print("  16. Thank you slide")
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(240, 248, 255)

thank_you_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1.5))
thank_you_frame = thank_you_box.text_frame
thank_you_frame.text = "Thank You"
thank_you_para = thank_you_frame.paragraphs[0]
thank_you_para.alignment = PP_ALIGN.CENTER
thank_you_para.font.size = Pt(54)
thank_you_para.font.bold = True
thank_you_para.font.color.rgb = TITLE_COLOR

subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.8))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "Questions?"
subtitle_para = subtitle_frame.paragraphs[0]
subtitle_para.alignment = PP_ALIGN.CENTER
subtitle_para.font.size = Pt(28)
subtitle_para.font.color.rgb = ACCENT_COLOR

# Save presentation
output_file = 'reports/Excel_Data_Analysis_Presentation_Final.pptx'
prs.save(output_file)

print("\n" + "=" * 70)
print("✓ PowerPoint Presentation Created!")
print("=" * 70)
print(f"\n📊 Presentation Details:")
print(f"  • Total Slides: 16")
print(f"  • File Size: {os.path.getsize(output_file) / 1024 / 1024:.2f} MB")
print(f"  • Location: {output_file}")
print("\n📁 Includes:")
print("  • Executive Summary")
print("  • Dataset Overview")
print("  • Active Leads Analysis")
print("  • Country & Regional Breakdown")
print("  • Email Bounced Analysis")
print("  • Industry & Lead Source Distribution")
print("  • Role Analysis (HR, IT, Finance, CEO, CFO)")
print("  • Data Quality Insights")
print("  • Key Findings & Recommendations")
print("=" * 70)

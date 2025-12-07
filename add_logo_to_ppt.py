"""
Add logo to presentation - First and Last slides
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import requests
from PIL import Image
import io

print("=" * 70)
print("Adding Logo to Presentation")
print("=" * 70)

# Download logo if needed
logo_path = 'reports/rheinicke_logo.png'

# Check if we need to create a placeholder logo
if not os.path.exists(logo_path):
    print("\nNote: Using placeholder. Please provide the actual logo file.")
    print("Place the logo at: reports/rheinicke_logo.png")

# Load existing presentation
ppt_file = 'reports/Excel_Data_Analysis_Presentation_Final.pptx'
if not os.path.exists(ppt_file):
    print(f"\n✗ Presentation file not found: {ppt_file}")
    exit(1)

print(f"\n✓ Loading presentation: {ppt_file}")
prs = Presentation(ppt_file)

# Define colors
TITLE_COLOR = RGBColor(31, 78, 121)
ACCENT_COLOR = RGBColor(68, 114, 196)
TEXT_COLOR = RGBColor(51, 51, 51)

print("\n📝 Updating First Slide...")
# Update first slide (index 0)
first_slide = prs.slides[0]

# Add logo to top right if logo exists
if os.path.exists(logo_path):
    try:
        first_slide.shapes.add_picture(logo_path, Inches(7.5), Inches(0.3), width=Inches(2))
        print("  ✓ Logo added to top right")
    except Exception as e:
        print(f"  ⚠ Could not add logo: {e}")
else:
    # Add text logo as placeholder
    logo_box = first_slide.shapes.add_textbox(Inches(7), Inches(0.3), Inches(2.5), Inches(0.8))
    logo_frame = logo_box.text_frame
    logo_frame.text = "Rheinicke\nIT Consulting"
    for paragraph in logo_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.RIGHT
        paragraph.font.size = Pt(12)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(139, 0, 0)
    print("  ⚠ Logo placeholder added (text only)")

# Add "Prepared by" text at bottom
prepared_box = first_slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.4))
prepared_frame = prepared_box.text_frame
prepared_frame.text = "Prepared by: Vijayapriya Settu"
prepared_para = prepared_frame.paragraphs[0]
prepared_para.alignment = PP_ALIGN.CENTER
prepared_para.font.size = Pt(14)
prepared_para.font.color.rgb = TEXT_COLOR
prepared_para.font.italic = True
print("  ✓ Added 'Prepared by: Vijayapriya Settu'")

print("\n📝 Updating Last Slide (Thank You)...")
# Update last slide
last_slide = prs.slides[-1]

# Add logo to top right
if os.path.exists(logo_path):
    try:
        last_slide.shapes.add_picture(logo_path, Inches(7.5), Inches(0.3), width=Inches(2))
        print("  ✓ Logo added to top right")
    except Exception as e:
        print(f"  ⚠ Could not add logo: {e}")
else:
    # Add text logo as placeholder
    logo_box = last_slide.shapes.add_textbox(Inches(7), Inches(0.3), Inches(2.5), Inches(0.8))
    logo_frame = logo_box.text_frame
    logo_frame.text = "Rheinicke\nIT Consulting"
    for paragraph in logo_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.RIGHT
        paragraph.font.size = Pt(12)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(139, 0, 0)
    print("  ⚠ Logo placeholder added (text only)")

# Add "Prepared by" text at bottom of last slide
prepared_box_last = last_slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.4))
prepared_frame_last = prepared_box_last.text_frame
prepared_frame_last.text = "Prepared by: Vijayapriya Settu"
prepared_para_last = prepared_frame_last.paragraphs[0]
prepared_para_last.alignment = PP_ALIGN.CENTER
prepared_para_last.font.size = Pt(12)
prepared_para_last.font.color.rgb = TEXT_COLOR
prepared_para_last.font.italic = True
print("  ✓ Added 'Prepared by: Vijayapriya Settu'")

# Save updated presentation
output_file = 'reports/Excel_Data_Analysis_Presentation_Final.pptx'
try:
    prs.save(output_file)
    print("\n" + "=" * 70)
    print("✓ Presentation Updated Successfully!")
    print("=" * 70)
    print(f"\n📊 Updated File: {output_file}")
    print(f"  • Logo added to first & last slides")
    print(f"  • Prepared by: Vijayapriya Settu")
    
    if not os.path.exists(logo_path):
        print(f"\n⚠ Note: To add the actual logo:")
        print(f"  1. Save the Rheinicke logo as: {logo_path}")
        print(f"  2. Run this script again")
    
    print("=" * 70)
except Exception as e:
    print(f"\n✗ Error saving presentation: {e}")
    # Try alternate filename
    output_file = 'reports/Excel_Data_Analysis_Presentation_Updated.pptx'
    prs.save(output_file)
    print(f"✓ Saved as: {output_file}")
